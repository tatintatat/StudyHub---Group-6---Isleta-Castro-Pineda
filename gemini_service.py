"""
gemini_service.py
─────────────────────────────────────────────────────────────────
StudyHub · Gemini AI Service  (REST API version — no SDK needed)
Calls the Gemini REST API directly using `requests`.
This avoids ALL SDK version conflicts and v1beta issues.
─────────────────────────────────────────────────────────────────
"""

from __future__ import annotations

import io
import json
import logging
import re
import textwrap
from dataclasses import dataclass, field
from enum import Enum
from typing import Any

import requests
import pdfplumber
from pptx import Presentation

logger = logging.getLogger(__name__)


# ══════════════════════════════════════════════════════════════
# Constants
# ══════════════════════════════════════════════════════════════

DEFAULT_MODEL     = "gemini-2.0-flash-lite"
MAX_CONTENT_CHARS = 120_000
GEMINI_API_BASE   = "https://generativelanguage.googleapis.com/v1/models"

# Models that work on the stable v1 API
AVAILABLE_MODELS = [
    "gemini-2.0-flash-lite",
    "gemini-2.0-flash",
    "gemini-1.5-flash",
    "gemini-1.5-pro",
]


# ══════════════════════════════════════════════════════════════
# Enums
# ══════════════════════════════════════════════════════════════

class GenType(str, Enum):
    FLASHCARD = "flashcard"
    QUIZ      = "quiz"
    REVIEWER  = "reviewer"


class QuizType(str, Enum):
    MCQ            = "mcq"
    TRUE_FALSE     = "truefalse"
    IDENTIFICATION = "identification"
    MIXED          = "mixed"


# ══════════════════════════════════════════════════════════════
# Result dataclass
# ══════════════════════════════════════════════════════════════

@dataclass
class AIResult:
    success:       bool
    items:         list[dict[str, Any]] = field(default_factory=list)
    raw:           str  = ""
    error:         str  = ""
    used_fallback: bool = False
    note:          str  = ""
    metadata:      dict[str, Any] = field(default_factory=dict)


# ══════════════════════════════════════════════════════════════
# File parsers
# ══════════════════════════════════════════════════════════════

class FileParser:
    @staticmethod
    def parse(file_bytes: bytes, filename: str) -> str:
        ext = filename.rsplit(".", 1)[-1].lower()
        if ext == "pdf":
            return FileParser._parse_pdf(file_bytes)
        if ext in ("pptx", "ppt"):
            return FileParser._parse_pptx(file_bytes)
        raise ValueError(f"Unsupported file type: .{ext}")

    @staticmethod
    def _parse_pdf(data: bytes) -> str:
        chunks: list[str] = []
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    chunks.append(text.strip())
                for table in page.extract_tables():
                    for row in table:
                        row_text = " | ".join(
                            cell.strip() if cell else ""
                            for cell in row if cell
                        )
                        if row_text:
                            chunks.append(row_text)
        return "\n\n".join(chunks)

    @staticmethod
    def _parse_pptx(data: bytes) -> str:
        prs = Presentation(io.BytesIO(data))
        chunks: list[str] = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            if slide.shapes.title and slide.shapes.title.text.strip():
                slide_texts.append(f"## {slide.shapes.title.text.strip()}")
            for shape in slide.shapes:
                if not shape.has_text_frame or shape == slide.shapes.title:
                    continue
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        slide_texts.append(line)
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Notes] {notes}")
            if slide_texts:
                chunks.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))
        return "\n\n".join(chunks)


# ══════════════════════════════════════════════════════════════
# Prompt builder
# ══════════════════════════════════════════════════════════════

class PromptBuilder:
    @staticmethod
    def flashcards(content: str, count: int) -> str:
        return textwrap.dedent(f"""
            You are an expert study-material creator.
            Generate exactly {count} high-quality flashcards from the document below.
            Rules:
            - Every flashcard must have "front" (question/term) and "back" (answer/definition).
            - Keep fronts concise. Keep backs 1-3 sentences max.
            - Do NOT repeat the same concept twice.
            - Output ONLY a valid JSON array, no markdown fences, no extra text.

            JSON format:
            [{{"front": "What is X?", "back": "X is ..."}}]

            DOCUMENT:
            \"\"\"
            {content[:MAX_CONTENT_CHARS]}
            \"\"\"
        """).strip()

    @staticmethod
    def quiz(content: str, count: int, quiz_type: QuizType) -> str:
        type_instructions = {
            QuizType.MCQ: '4 options labelled a) b) c) d). "type":"mcq"',
            QuizType.TRUE_FALSE: 'options:["True","False"]. "type":"truefalse"',
            QuizType.IDENTIFICATION: 'options:[]. Short phrase answer. "type":"identification"',
            QuizType.MIXED: 'Mix of mcq, truefalse, identification (~1/3 each).',
        }
        return textwrap.dedent(f"""
            You are an expert educator. Generate exactly {count} quiz questions.
            Type: {quiz_type.value}. Instructions: {type_instructions[quiz_type]}
            Base questions ONLY on the document. Output ONLY a valid JSON array, no extra text.

            JSON format:
            [{{"question":"...","type":"mcq","options":["a) ...","b) ...","c) ...","d) ..."],"answer":"a) ...","explanation":"..."}}]

            DOCUMENT:
            \"\"\"
            {content[:MAX_CONTENT_CHARS]}
            \"\"\"
        """).strip()

    @staticmethod
    def reviewer(content: str) -> str:
        return textwrap.dedent(f"""
            You are an expert educator. Create a structured JSON study reviewer.
            Output ONLY a valid JSON object, no markdown fences, no extra text.

            Format:
            {{"title":"...","summary":"...","key_concepts":[{{"term":"...","definition":"...","example":"..."}}],
            "important_facts":["..."],"common_misconceptions":["..."],
            "key_formulas_or_rules":["..."],"study_tips":["..."],"possible_exam_topics":["..."]}}

            DOCUMENT:
            \"\"\"
            {content[:MAX_CONTENT_CHARS]}
            \"\"\"
        """).strip()


# ══════════════════════════════════════════════════════════════
# Response parser
# ══════════════════════════════════════════════════════════════

class ResponseParser:
    @staticmethod
    def parse_array(raw: str) -> list[dict]:
        cleaned = re.sub(r"```(?:json)?", "", raw, flags=re.IGNORECASE).strip().rstrip("`")
        match = re.search(r"\[[\s\S]*\]", cleaned)
        if not match:
            raise ValueError("No JSON array found in response")
        return json.loads(match.group(0))

    @staticmethod
    def parse_object(raw: str) -> dict:
        cleaned = re.sub(r"```(?:json)?", "", raw, flags=re.IGNORECASE).strip().rstrip("`")
        match = re.search(r"\{[\s\S]*\}", cleaned)
        if not match:
            raise ValueError("No JSON object found in response")
        return json.loads(match.group(0))

    @staticmethod
    def validate_flashcards(items: list[dict]) -> list[dict]:
        return [
            {"front": str(i["front"]).strip(), "back": str(i["back"]).strip()}
            for i in items
            if isinstance(i, dict) and i.get("front") and i.get("back")
        ]

    @staticmethod
    def validate_quiz(items: list[dict]) -> list[dict]:
        valid = []
        for i in items:
            if not isinstance(i, dict) or not i.get("question") or not i.get("answer"):
                continue
            q_type = str(i.get("type", "mcq")).lower()
            if q_type not in {"mcq", "truefalse", "identification"}:
                q_type = "mcq"
            valid.append({
                "question":    str(i["question"]).strip(),
                "type":        q_type,
                "options":     list(i.get("options", [])),
                "answer":      str(i["answer"]).strip(),
                "explanation": str(i.get("explanation", "")).strip(),
            })
        return valid


# ══════════════════════════════════════════════════════════════
# Main service — uses REST API directly (no SDK)
# ══════════════════════════════════════════════════════════════

class GeminiService:
    """
    Calls Gemini REST API directly via requests.
    Uses stable v1 endpoint — no SDK version conflicts.
    """

    def __init__(self, api_key: str, model: str = DEFAULT_MODEL) -> None:
        if not api_key:
            raise ValueError("Gemini API key is required.")
        self._api_key    = api_key
        self._model_name = model
        logger.info("GeminiService ready (REST) model=%s", model)

    # ── Public API ────────────────────────────────────────────

    def generate_from_file(self, file_bytes, filename, gen_type="flashcard",
                           quiz_type="mcq", count=10) -> AIResult:
        try:
            text = FileParser.parse(file_bytes, filename)
        except Exception as exc:
            return AIResult(success=False, error=f"Could not read file: {exc}")
        if not text.strip():
            return AIResult(success=False, error="Document appears empty or image-only.")
        return self._generate(text, gen_type, quiz_type, count)

    def generate_from_text(self, text, gen_type="flashcard",
                           quiz_type="mcq", count=10) -> AIResult:
        if not text.strip():
            return AIResult(success=False, error="No text provided.")
        return self._generate(text, gen_type, quiz_type, count)

    def generate_flashcards(self, text, count=10) -> AIResult:
        return self._generate(text, GenType.FLASHCARD, QuizType.MCQ, count)

    def generate_quiz(self, text, quiz_type="mcq", count=10) -> AIResult:
        return self._generate(text, GenType.QUIZ, quiz_type, count)

    def generate_reviewer(self, text) -> AIResult:
        return self._generate(text, GenType.REVIEWER, QuizType.MCQ, 0)

    def extract_text_from_file(self, file_bytes, filename) -> str:
        return FileParser.parse(file_bytes, filename)

    # ── Internal ──────────────────────────────────────────────

    def _generate(self, text, gen_type, quiz_type, count) -> AIResult:
        try:
            g = GenType(gen_type)
            q = QuizType(quiz_type) if gen_type == "quiz" else QuizType.MCQ
        except ValueError:
            return AIResult(success=False, error=f"Unknown gen_type '{gen_type}'.")

        if g == GenType.FLASHCARD:
            prompt = PromptBuilder.flashcards(text, count)
        elif g == GenType.QUIZ:
            prompt = PromptBuilder.quiz(text, count, q)
        else:
            prompt = PromptBuilder.reviewer(text)

        try:
            raw = self._call_rest(prompt)
        except Exception as exc:
            logger.error("Gemini API error: %s", exc)
            return AIResult(success=False, error=f"Gemini API error: {exc}")

        return self._parse_response(raw, g, count)

    def _call_rest(self, prompt: str) -> str:
        """Call Gemini stable v1 REST API directly."""
        url = f"{GEMINI_API_BASE}/{self._model_name}:generateContent?key={self._api_key}"
        payload = {
            "contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {
                "temperature": 0.7,
                "topP": 0.95,
                "maxOutputTokens": 8192,
            },
        }
        resp = requests.post(url, json=payload, timeout=120)

        if resp.status_code == 429:
            raise Exception(f"429 RESOURCE_EXHAUSTED. {resp.json()}")
        if resp.status_code == 404:
            raise Exception(f"404 NOT_FOUND — model '{self._model_name}' unavailable. {resp.json()}")
        if not resp.ok:
            raise Exception(f"{resp.status_code} {resp.text}")

        data = resp.json()
        try:
            return data["candidates"][0]["content"]["parts"][0]["text"]
        except (KeyError, IndexError):
            raise RuntimeError(f"Unexpected response format: {data}")

    def _parse_response(self, raw: str, g: GenType, requested_count: int) -> AIResult:
        try:
            if g == GenType.FLASHCARD:
                items = ResponseParser.validate_flashcards(ResponseParser.parse_array(raw))
                if not items:
                    raise ValueError("No valid flashcards parsed.")
                return AIResult(success=True, items=items, raw=raw,
                                metadata={"count": len(items), "requested": requested_count})

            if g == GenType.QUIZ:
                items = ResponseParser.validate_quiz(ResponseParser.parse_array(raw))
                if not items:
                    raise ValueError("No valid quiz questions parsed.")
                return AIResult(success=True, items=items, raw=raw,
                                metadata={"count": len(items), "requested": requested_count})

            if g == GenType.REVIEWER:
                obj = ResponseParser.parse_object(raw)
                return AIResult(success=True, items=[obj], raw=raw,
                                metadata={"type": "reviewer"})

        except (json.JSONDecodeError, ValueError) as exc:
            logger.warning("Parse error: %s — falling back to raw.", exc)
            return AIResult(success=True, items=[], raw=raw, used_fallback=True,
                            note="AI responded but output could not be auto-parsed.")
        except Exception as exc:
            return AIResult(success=False, error=str(exc), raw=raw)

        return AIResult(success=False, error="Unexpected generation type.")
